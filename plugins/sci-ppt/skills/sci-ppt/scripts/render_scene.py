#!/usr/bin/env python3
"""Render Sci-PPT scene JSON into editable native PowerPoint objects.

The renderer intentionally favors ordinary PowerPoint primitives so users can
move, recolor, resize, relabel, and delete parts after generation.
"""

import argparse
import json
import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def rgb(value):
    if value is None:
        return None
    value = str(value).lstrip("#")
    if len(value) != 6:
        raise ValueError(f"Expected #RRGGBB, got {value!r}")
    return RGBColor.from_string(value.upper())


def apply_style(shape, obj):
    fill = obj.get("fill", "#FFFFFF")
    if hasattr(shape, "fill"):
        if fill is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(fill)

    if hasattr(shape, "line"):
        line = obj.get("line", "#333333")
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = rgb(line)
            shape.line.width = Pt(float(obj.get("line_width", 1.25)))

    if "rotation" in obj:
        shape.rotation = float(obj["rotation"])


def add_text(slide, obj):
    shape = slide.shapes.add_textbox(
        Inches(float(obj["x"])), Inches(float(obj["y"])),
        Inches(float(obj["w"])), Inches(float(obj["h"]))
    )
    shape.rotation = float(obj.get("rotation", 0))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = bool(obj.get("word_wrap", True))
    tf.margin_left = Inches(float(obj.get("margin", 0.03)))
    tf.margin_right = Inches(float(obj.get("margin", 0.03)))
    tf.margin_top = Inches(float(obj.get("margin", 0.02)))
    tf.margin_bottom = Inches(float(obj.get("margin", 0.02)))
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(str(obj.get("valign", "middle")).lower(), MSO_ANCHOR.MIDDLE)

    lines = str(obj.get("text", "")).split("\n")
    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }.get(str(obj.get("align", "left")).lower(), PP_ALIGN.LEFT)
        p.space_after = Pt(float(obj.get("space_after", 0)))
        if p.runs:
            run = p.runs[0]
            run.font.size = Pt(float(obj.get("font_size", 18)))
            run.font.name = obj.get("font", "Arial")
            run.font.bold = bool(obj.get("bold", False))
            run.font.italic = bool(obj.get("italic", False))
            run.font.color.rgb = rgb(obj.get("color", "#222222"))
    return shape


def add_polygon(slide, obj):
    points = obj["points"]
    if len(points) < 3:
        raise ValueError("polygon needs at least 3 points")
    builder = slide.shapes.build_freeform(Inches(points[0][0]), Inches(points[0][1]))
    segments = [(Inches(x), Inches(y)) for x, y in points[1:]]
    builder.add_line_segments(segments, close=True)
    shape = builder.convert_to_shape()
    apply_style(shape, obj)
    return shape


def add_connector(slide, obj, arrow=False):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(float(obj["x1"])), Inches(float(obj["y1"])),
        Inches(float(obj["x2"])), Inches(float(obj["y2"]))
    )
    shape.line.color.rgb = rgb(obj.get("line", "#333333"))
    shape.line.width = Pt(float(obj.get("line_width", 1.5)))
    if arrow:
        ln = shape.line._get_or_add_ln()
        tail = ln.makeelement(f"{{{A_NS}}}tailEnd")
        tail.set("type", obj.get("arrowhead", "triangle"))
        tail.set("w", obj.get("arrow_width", "med"))
        tail.set("len", obj.get("arrow_length", "med"))
        ln.append(tail)
    return shape


def add_basic_shape(slide, obj, shape_type):
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(float(obj["x"])), Inches(float(obj["y"])),
        Inches(float(obj["w"])), Inches(float(obj["h"]))
    )
    apply_style(shape, obj)
    return shape


def add_label(slide, text, x, y, w, h, **style):
    obj = {
        "type": "text", "text": text, "x": x, "y": y, "w": w, "h": h,
        "font_size": style.get("font_size", 11), "font": style.get("font", "Arial"),
        "color": style.get("color", "#222222"), "align": style.get("align", "center"),
        "valign": "middle", "bold": style.get("bold", False),
    }
    return add_text(slide, obj)


def add_particle_cluster(slide, obj):
    """Draw a cluster as independent editable circles."""
    cx, cy = float(obj["x"]), float(obj["y"])
    w, h = float(obj["w"]), float(obj["h"])
    count = int(obj.get("count", 9))
    diameter = float(obj.get("diameter", min(w, h) / 4.5))
    fill = obj.get("fill", "#F4B183")
    line = obj.get("line", "#B45F06")
    for i in range(count):
        angle = i * 2.399963229728653
        radius = min(w, h) * 0.34 * math.sqrt((i + 0.5) / max(count, 1))
        px = cx + w / 2 + radius * math.cos(angle) - diameter / 2
        py = cy + h / 2 + radius * math.sin(angle) - diameter / 2
        add_basic_shape(slide, {
            "x": px, "y": py, "w": diameter, "h": diameter,
            "fill": fill, "line": line, "line_width": obj.get("line_width", 0.8)
        }, MSO_SHAPE.OVAL)
    if obj.get("label"):
        add_label(slide, obj["label"], cx, cy + h + 0.03, w, 0.28,
                  font_size=obj.get("font_size", 10))


def add_layered_block(slide, obj):
    """Draw substrate/coating stacks as separate editable rectangles."""
    x, y, w, h = map(float, (obj["x"], obj["y"], obj["w"], obj["h"]))
    layers = obj.get("layers", [])
    if not layers:
        layers = [{"fill": "#D9EAD3", "label": "substrate"}]
    total = sum(float(layer.get("ratio", 1)) for layer in layers)
    cursor = y + h
    for layer in reversed(layers):
        lh = h * float(layer.get("ratio", 1)) / total
        cursor -= lh
        add_basic_shape(slide, {
            "x": x, "y": cursor, "w": w, "h": lh,
            "fill": layer.get("fill", "#D9EAD3"),
            "line": layer.get("line", obj.get("line", "#666666")),
            "line_width": layer.get("line_width", obj.get("line_width", 0.8)),
        }, MSO_SHAPE.RECTANGLE)
        if layer.get("label"):
            add_label(slide, layer["label"], x, cursor, w, lh,
                      font_size=layer.get("font_size", obj.get("font_size", 9)),
                      color=layer.get("text_color", "#222222"))


def add_membrane(slide, obj):
    """Draw a stylized lipid-like bilayer using editable circles and lines."""
    x, y, w = float(obj["x"]), float(obj["y"]), float(obj["w"])
    spacing = float(obj.get("spacing", 0.28))
    head = float(obj.get("head", 0.12))
    tail = float(obj.get("tail", 0.25))
    count = max(2, int(w / spacing))
    for row, direction in ((0, 1), (1, -1)):
        yy = y + row * (head + tail)
        for i in range(count):
            xx = x + i * w / max(count - 1, 1) - head / 2
            add_basic_shape(slide, {
                "x": xx, "y": yy, "w": head, "h": head,
                "fill": obj.get("head_fill", "#6FA8DC"),
                "line": obj.get("line", "#3D85C6"), "line_width": 0.6,
            }, MSO_SHAPE.OVAL)
            cx = xx + head / 2
            y1 = yy + (head if direction == 1 else 0)
            y2 = y1 + direction * tail
            add_connector(slide, {
                "x1": cx, "y1": y1, "x2": cx - 0.035, "y2": y2,
                "line": obj.get("tail_line", "#666666"), "line_width": 0.6,
            })
            add_connector(slide, {
                "x1": cx, "y1": y1, "x2": cx + 0.035, "y2": y2,
                "line": obj.get("tail_line", "#666666"), "line_width": 0.6,
            })


def add_cell(slide, obj):
    """Draw a simple editable cell: body, nucleus, optional particles."""
    x, y, w, h = map(float, (obj["x"], obj["y"], obj["w"], obj["h"]))
    add_basic_shape(slide, {
        "x": x, "y": y, "w": w, "h": h,
        "fill": obj.get("fill", "#D9EAD3"), "line": obj.get("line", "#6AA84F"),
        "line_width": obj.get("line_width", 1.5),
    }, MSO_SHAPE.OVAL)
    nw, nh = w * 0.36, h * 0.36
    add_basic_shape(slide, {
        "x": x + w * 0.34, "y": y + h * 0.32, "w": nw, "h": nh,
        "fill": obj.get("nucleus_fill", "#C9B2E5"),
        "line": obj.get("nucleus_line", "#674EA7"), "line_width": 1.0,
    }, MSO_SHAPE.OVAL)
    particles = int(obj.get("particles", 0))
    for i in range(particles):
        a = (i + 1) * 2.2
        px = x + w * (0.5 + 0.33 * math.cos(a))
        py = y + h * (0.5 + 0.33 * math.sin(a))
        d = min(w, h) * 0.07
        add_basic_shape(slide, {
            "x": px - d / 2, "y": py - d / 2, "w": d, "h": d,
            "fill": obj.get("particle_fill", "#F6B26B"),
            "line": obj.get("particle_line", "#E69138"), "line_width": 0.5,
        }, MSO_SHAPE.OVAL)
    if obj.get("label"):
        add_label(slide, obj["label"], x, y + h + 0.03, w, 0.3,
                  font_size=obj.get("font_size", 10))


def add_beaker(slide, obj):
    """Draw an editable simplified beaker with optional liquid and label."""
    x, y, w, h = map(float, (obj["x"], obj["y"], obj["w"], obj["h"]))
    # Vessel outline is a trapezoid-like polygon.
    points = [
        [x + w * 0.18, y], [x + w * 0.82, y],
        [x + w * 0.92, y + h], [x + w * 0.08, y + h],
    ]
    add_polygon(slide, {
        "points": points, "fill": obj.get("fill", "#FFFFFF"),
        "line": obj.get("line", "#666666"), "line_width": obj.get("line_width", 1.2),
    })
    liquid_fraction = max(0.0, min(1.0, float(obj.get("liquid_fraction", 0.55))))
    lh = h * liquid_fraction * 0.86
    add_basic_shape(slide, {
        "x": x + w * 0.13, "y": y + h - lh - h * 0.03,
        "w": w * 0.74, "h": lh,
        "fill": obj.get("liquid_fill", "#CFE2F3"), "line": None,
    }, MSO_SHAPE.RECTANGLE)
    if obj.get("label"):
        add_label(slide, obj["label"], x + w * 0.1, y + h * 0.48, w * 0.8, h * 0.25,
                  font_size=obj.get("font_size", 10))


def add_droplet(slide, obj):
    add_basic_shape(slide, obj, MSO_SHAPE.TEAR)


def render_object(slide, obj):
    kind = obj["type"]
    if kind == "text":
        return add_text(slide, obj)
    if kind == "polygon":
        return add_polygon(slide, obj)
    if kind == "line":
        return add_connector(slide, obj, arrow=False)
    if kind == "arrow":
        return add_connector(slide, obj, arrow=True)
    if kind == "particle_cluster":
        return add_particle_cluster(slide, obj)
    if kind == "layered_block":
        return add_layered_block(slide, obj)
    if kind == "membrane":
        return add_membrane(slide, obj)
    if kind == "cell":
        return add_cell(slide, obj)
    if kind == "beaker":
        return add_beaker(slide, obj)
    if kind == "droplet":
        return add_droplet(slide, obj)

    shape_type = {
        "rect": MSO_SHAPE.RECTANGLE,
        "round_rect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "ellipse": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "diamond": MSO_SHAPE.DIAMOND,
        "hexagon": MSO_SHAPE.HEXAGON,
        "chevron": MSO_SHAPE.CHEVRON,
        "cylinder": MSO_SHAPE.CAN,
        "cloud": MSO_SHAPE.CLOUD,
        "star": MSO_SHAPE.STAR_5_POINT,
    }.get(kind)
    if shape_type is None:
        raise ValueError(f"Unsupported object type: {kind}")
    return add_basic_shape(slide, obj, shape_type)


def validate_scene(scene):
    if not isinstance(scene, dict):
        raise ValueError("Scene must be a JSON object")
    if "objects" not in scene or not isinstance(scene["objects"], list):
        raise ValueError("Scene requires an objects array")
    for i, obj in enumerate(scene["objects"]):
        if not isinstance(obj, dict) or "type" not in obj:
            raise ValueError(f"Object {i} must be an object with a type")


def render(scene, output):
    validate_scene(scene)
    prs = Presentation()
    slide_cfg = scene.get("slide", {})
    prs.slide_width = Inches(float(slide_cfg.get("width", 13.333)))
    prs.slide_height = Inches(float(slide_cfg.get("height", 7.5)))
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(scene.get("background", "#FFFFFF"))

    for obj in scene.get("objects", []):
        render_object(slide, obj)

    output = Path(output)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)

    check = Presentation(output)
    if len(check.slides) < 1 or len(check.slides[0].shapes) < 1:
        raise RuntimeError("PPTX verification failed: no editable shapes found")
    return output


def main():
    ap = argparse.ArgumentParser(description="Render a Sci-PPT JSON scene into editable PowerPoint objects.")
    ap.add_argument("--scene", required=True)
    ap.add_argument("--output", required=True)
    args = ap.parse_args()
    with open(args.scene, "r", encoding="utf-8") as f:
        scene = json.load(f)
    out = render(scene, args.output)
    print(out)


if __name__ == "__main__":
    main()
