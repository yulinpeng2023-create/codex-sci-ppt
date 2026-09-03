#!/usr/bin/env python3
"""Render Sci-PPT geometry cache into editable native PowerPoint objects.

The custom-geometry OOXML approach is adapted from the MIT-licensed Cell_ppt
renderer. See THIRD_PARTY_NOTICES.md.
"""
from __future__ import annotations

import argparse
import json
import os
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Emu, Pt

EMU_PER_PT = 12700.0
PATH_EXTENT = 100000


def esc(value):
    return (
        str(value).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        .replace('"', "&quot;").replace("'", "&apos;")
    )


def same(first, second):
    return abs(first[0] - second[0]) < 1e-6 and abs(first[1] - second[1]) < 1e-6


def rgb_hex(value, fallback="000000"):
    if not value:
        return fallback
    return "".join(f"{max(0, min(255, int(component))):02X}" for component in value[:3])


def alpha_xml(opacity):
    alpha = max(0, min(100000, round(float(opacity) * 1000)))
    return f'<a:alpha val="{alpha}"/>' if alpha < 100000 else ""


def map_point(point, transform):
    scale, offset_x, offset_y, view_x, view_y = transform
    return (
        offset_x + (point[0] - view_x) * scale,
        offset_y + (point[1] - view_y) * scale,
    )


def add_freeform(slide, subpath, paint, name, shape_id, transform):
    mapped = [
        {key: map_point(point[key], transform) for key in ("a", "l", "r")}
        for point in subpath["points"]
    ]
    coordinates = [point[key] for point in mapped for key in ("a", "l", "r")]
    min_x = min(point[0] for point in coordinates)
    min_y = min(point[1] for point in coordinates)
    max_x = max(point[0] for point in coordinates)
    max_y = max(point[1] for point in coordinates)
    width = max(max_x - min_x, 1.0 / EMU_PER_PT)
    height = max(max_y - min_y, 1.0 / EMU_PER_PT)

    def local(point):
        return (
            round((point[0] - min_x) / width * PATH_EXTENT),
            round((point[1] - min_y) / height * PATH_EXTENT),
        )

    first_x, first_y = local(mapped[0]["a"])
    commands = [f'<a:moveTo><a:pt x="{first_x}" y="{first_y}"/></a:moveTo>']
    for index in range(1, len(mapped)):
        previous = mapped[index - 1]
        current = mapped[index]
        end_x, end_y = local(current["a"])
        if same(previous["r"], previous["a"]) and same(current["l"], current["a"]):
            commands.append(f'<a:lnTo><a:pt x="{end_x}" y="{end_y}"/></a:lnTo>')
        else:
            c1x, c1y = local(previous["r"])
            c2x, c2y = local(current["l"])
            commands.append(
                f'<a:cubicBezTo><a:pt x="{c1x}" y="{c1y}"/>'
                f'<a:pt x="{c2x}" y="{c2y}"/><a:pt x="{end_x}" y="{end_y}"/></a:cubicBezTo>'
            )
    if subpath.get("closed"):
        commands.append("<a:close/>")

    fill = "<a:noFill/>"
    if paint.get("filled") and subpath.get("closed"):
        fill = (
            f'<a:solidFill><a:srgbClr val="{rgb_hex(paint.get("fillColor"))}">'
            f'{alpha_xml(paint.get("opacity", 100))}</a:srgbClr></a:solidFill>'
        )

    line = '<a:ln><a:noFill/></a:ln>'
    if paint.get("stroked"):
        weight = max(3175, round(float(paint.get("strokeWidth", 1)) * transform[0] * EMU_PER_PT))
        line = (
            f'<a:ln w="{weight}"><a:solidFill><a:srgbClr val="{rgb_hex(paint.get("strokeColor"))}">'
            f'{alpha_xml(paint.get("opacity", 100))}</a:srgbClr></a:solidFill></a:ln>'
        )

    xml = f'''<p:sp {nsdecls("a", "p")}>
      <p:nvSpPr><p:cNvPr id="{shape_id}" name="{esc(name)}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{round(min_x * EMU_PER_PT)}" y="{round(min_y * EMU_PER_PT)}"/><a:ext cx="{max(1, round(width * EMU_PER_PT))}" cy="{max(1, round(height * EMU_PER_PT))}"/></a:xfrm>
        <a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/><a:rect l="l" t="t" r="r" b="b"/>
          <a:pathLst><a:path w="{PATH_EXTENT}" h="{PATH_EXTENT}">{''.join(commands)}</a:path></a:pathLst>
        </a:custGeom>{fill}{line}
      </p:spPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
    </p:sp>'''
    slide.shapes._spTree.insert_element_before(parse_xml(xml), "p:extLst")


def add_text(slide, atom, transform):
    text = atom["text"]
    x, y = map_point(text["position"], transform)
    size = max(4.0, float(text["fontSize"]) * transform[0])
    width = max(size * 2.0, size * len(str(text["contents"])) * 0.7)
    height = size * 1.5
    anchor = str(text.get("textAnchor", "start")).lower()
    if anchor == "middle":
        x -= width / 2.0
    elif anchor == "end":
        x -= width

    shape = slide.shapes.add_textbox(
        Emu(round(x * EMU_PER_PT)),
        Emu(round((y - size * 1.05) * EMU_PER_PT)),
        Emu(round(width * EMU_PER_PT)),
        Emu(round(height * EMU_PER_PT)),
    )
    shape.name = atom["objectName"]
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = str(text["contents"])
    font = run.font
    font.name = text.get("fontFamily", "Arial")
    font.size = Pt(size)
    weight = str(text.get("fontWeight", "")).lower()
    font.bold = weight == "bold" or (weight.isdigit() and int(weight) >= 600)
    font.italic = str(text.get("fontStyle", "")).lower() == "italic"
    font.color.rgb = RGBColor(*[int(value) for value in text.get("fillColor", [0, 0, 0])])
    shape.rotation = float(text.get("rotationDegrees", 0))


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--geometry-cache", required=True, type=Path)
    parser.add_argument("--output-pptx", required=True, type=Path)
    parser.add_argument("--input-pptx", type=Path)
    parser.add_argument("--slide-index", type=int, default=0)
    args = parser.parse_args()

    cache = json.loads(args.geometry_cache.read_text(encoding="utf-8"))
    if int(cache.get("schema_version", 0)) != 3:
        raise SystemExit("Unsupported geometry cache schema")

    prs = Presentation(str(args.input_pptx)) if args.input_pptx else Presentation()
    if args.slide_index:
        if not 1 <= args.slide_index <= len(prs.slides):
            raise SystemExit("slide-index is out of range")
        slide = prs.slides[args.slide_index - 1]
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width / EMU_PER_PT
    slide_height = prs.slide_height / EMU_PER_PT
    view_x, view_y, view_width, view_height = map(float, cache["view_box"])
    margin = 18.0
    scale = min(
        (slide_width - 2 * margin) / view_width,
        (slide_height - 2 * margin) / view_height,
    )
    transform = (
        scale,
        (slide_width - view_width * scale) / 2.0,
        (slide_height - view_height * scale) / 2.0,
        view_x,
        view_y,
    )

    next_id = max((shape.shape_id for shape in slide.shapes), default=1) + 1
    created = 0
    for batch in cache["batches"]:
        for atom_index in batch["atom_indices"]:
            atom = cache["atoms"][int(atom_index)]
            if atom["kind"] == "text":
                add_text(slide, atom, transform)
                next_id += 1
                created += 1
                continue
            paints = atom.get("paintParts") or [{}]
            for part_index, subpath in enumerate(atom.get("subpaths", [])):
                if len(subpath.get("points", [])) < 2:
                    continue
                add_freeform(
                    slide,
                    subpath,
                    paints[min(part_index, len(paints) - 1)],
                    f'{atom["objectName"]}_PART_{part_index:03d}',
                    next_id,
                    transform,
                )
                next_id += 1
                created += 1

    args.output_pptx.parent.mkdir(parents=True, exist_ok=True)
    handle, temporary_name = tempfile.mkstemp(
        prefix="._scippt-",
        suffix=".pptx",
        dir=args.output_pptx.parent,
    )
    os.close(handle)
    temporary = Path(temporary_name)
    try:
        prs.save(temporary)
        Presentation(temporary)
        os.replace(temporary, args.output_pptx)
    finally:
        if temporary.exists():
            temporary.unlink()

    reopened = Presentation(args.output_pptx)
    print(
        json.dumps(
            {
                "ok": True,
                "backend": "editable-ooxml",
                "output_pptx": str(args.output_pptx.resolve()),
                "native_object_count": created,
                "slide_count": len(reopened.slides),
            },
            ensure_ascii=False,
        )
    )


if __name__ == "__main__":
    main()
