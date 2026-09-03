#!/usr/bin/env python3
"""Local raster-to-editable-PPT tracing for simple scientific diagrams."""
import argparse
from pathlib import Path
import cv2
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches


def trace(image_path, output_path, max_shapes=400):
    image = cv2.imread(str(image_path))
    if image is None:
        raise FileNotFoundError(image_path)
    h, w = image.shape[:2]
    pixels = image.reshape((-1, 3)).astype(np.float32)
    criteria = (cv2.TERM_CRITERIA_EPS + cv2.TERM_CRITERIA_MAX_ITER, 25, 1.0)
    _, labels, centers = cv2.kmeans(pixels, 8, None, criteria, 3, cv2.KMEANS_PP_CENTERS)
    labels = labels.reshape(h, w)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(13.333 * h / w)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sx = 13.333 / w
    sy = (13.333 * h / w) / h
    minimum = w * h * 0.0002
    found = []

    for index in range(len(centers)):
        mask = (labels == index).astype(np.uint8) * 255
        contours, _ = cv2.findContours(mask, cv2.RETR_LIST, cv2.CHAIN_APPROX_SIMPLE)
        for contour in contours:
            area = cv2.contourArea(contour)
            if minimum <= area <= w * h * 0.95:
                found.append((area, index, contour))

    found.sort(reverse=True, key=lambda item: item[0])
    for _, index, contour in found[:max_shapes]:
        perimeter = cv2.arcLength(contour, True)
        approx = cv2.approxPolyDP(contour, 0.005 * perimeter, True).reshape(-1, 2)
        if len(approx) < 3:
            continue
        first_x, first_y = approx[0]
        builder = slide.shapes.build_freeform(Inches(first_x * sx), Inches(first_y * sy))
        for x, y in approx[1:]:
            builder.add_line_segments([(Inches(x * sx), Inches(y * sy))], close=False)
        shape = builder.convert_to_shape()
        b, g, r = centers[index].astype(int)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(int(r), int(g), int(b))
        shape.line.fill.background()

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    check = Presentation(output)
    count = len(check.slides[0].shapes)
    if count == 0:
        raise RuntimeError('No editable shapes were produced.')
    print(f'{output} ({count} editable shapes)')


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--input-image', required=True)
    parser.add_argument('--output', required=True)
    parser.add_argument('--max-shapes', type=int, default=400)
    args = parser.parse_args()
    trace(args.input_image, args.output, args.max_shapes)


if __name__ == '__main__':
    main()
