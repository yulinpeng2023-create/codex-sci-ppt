#!/usr/bin/env python3
from __future__ import annotations

import json
import tempfile
from pathlib import Path

from pptx import Presentation

from render_bamboo_template import render


def main():
    here = Path(__file__).resolve().parent
    config_path = here.parent / "templates" / "bamboo_cross_section.json"
    config = json.loads(config_path.read_text(encoding="utf-8"))

    with tempfile.TemporaryDirectory(prefix="codex-sci-ppt-bamboo-") as tmp:
        output = Path(tmp) / "bamboo-template.pptx"
        out, shape_count = render(config, output)
        prs = Presentation(out)
        assert len(prs.slides) == 1
        assert shape_count >= 380
        assert len(prs.slides[0].shapes) == shape_count
        assert prs.slides[0].shapes[0].shape_type is not None
        print(json.dumps({"ok": True, "shapes": shape_count, "template": str(config_path)}, ensure_ascii=False))


if __name__ == "__main__":
    main()
