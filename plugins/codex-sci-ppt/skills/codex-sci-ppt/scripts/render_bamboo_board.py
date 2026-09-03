#!/usr/bin/env python3
"""Render a reusable editable bamboo board in 5 x 2 x 0.5 proportions.

The physical proportions are length:width:thickness = 5:2:0.5 (10:4:1).
The long direction is bamboo longitudinal direction, so vascular bundles are
shown on the transverse end face (2 x 0.5), while top and long side faces use
longitudinal grain cues.
"""
from __future__ import annotations

import argparse
import json
import random
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#").upper())


def add_polygon(slide, points, fill, line="#A66A00", line_width=0.8):
    builder = slide.shapes.build_freeform(Inches(points[0][0]), Inches(points[0][1]))
    builder.add_line_segments([(Inches(x), Inches(y)) for x, y in points[1:]], close=True)
    shape = builder.convert_to_shape()
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(float(line_width))
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
    return shape


def add_oval(slide, x, y, w, h, fill, line=None, line_width=0.4):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
    return shape


def add_vascular_bundle(slide, cx, cy, size, dark="#6C3B05", gold="#DDAE3E", rotation=0.0):
    """Small editable stylized bamboo vascular bundle."""
    for dx in (-0.29, 0.29):
        s = add_oval(slide, cx + dx*size - 0.14*size, cy - 0.11*size,
                     0.28*size, 0.28*size, dark)
        s.rotation = rotation
        g = add_oval(slide, cx + dx*size - 0.085*size, cy - 0.055*size,
                     0.17*size, 0.17*size, gold)
        g.rotation = rotation
    s = add_oval(slide, cx - 0.16*size, cy - 0.34*size, 0.32*size, 0.26*size, dark)
    s.rotation = rotation
    s = add_oval(slide, cx - 0.15*size, cy - 0.09*size, 0.30*size, 0.25*size, gold)
    s.rotation = rotation
    s = slide.shapes.add_shape(MSO_SHAPE.TEAR, Inches(cx - 0.13*size), Inches(cy + 0.08*size), Inches(0.26*size), Inches(0.34*size))
    s.rotation = 180 + rotation
    s.fill.solid(); s.fill.fore_color.rgb = rgb(dark); s.line.fill.background()
    for dx in (-0.045, 0.045):
        add_oval(slide, cx + dx*size - 0.012*size, cy - 0.002*size,
                 0.024*size, 0.024*size, dark)


def projected_board_geometry(cfg):
    x = float(cfg.get("x", 1.55))
    y = float(cfg.get("y", 3.0))
    scale = float(cfg.get("scale", 1.0))
    # Display dimensions follow physical 10:4:1 proportions while using a 3/4 projection.
    length = 7.0 * scale
    width_dx = 2.8 * scale
    width_dy = -1.0 * scale
    thickness = 0.70 * scale
    p0 = (x, y)
    p1 = (x + length, y)
    p2 = (x + length + width_dx, y + width_dy)
    p3 = (x + width_dx, y + width_dy)
    top = [p0, p1, p2, p3]
    end = [p0, p3, (p3[0], p3[1] + thickness), (p0[0], p0[1] + thickness)]
    side = [p1, p2, (p2[0], p2[1] + thickness), (p1[0], p1[1] + thickness)]
    return top, end, side


def sample_nonuniform_bundles(rng, end_poly, count=18):
    """Irregular distribution with radial density/size gradient."""
    p0, p1, p2, p3 = end_poly
    accepted = []
    attempts = 0
    while len(accepted) < count and attempts < 6000:
        attempts += 1
        # Mixture distribution: locally denser zones plus sparse gaps; no grid.
        if rng.random() < 0.62:
            u = min(0.98, max(0.02, rng.betavariate(1.15, 1.15)))
        else:
            center = rng.choice((0.16, 0.38, 0.68, 0.86))
            u = min(0.98, max(0.02, rng.gauss(center, 0.075)))
        # Slightly denser toward the outer side (v near 0).
        v = rng.random() ** 1.45
        tx = p0[0] + (p1[0]-p0[0]) * u
        ty = p0[1] + (p1[1]-p0[1]) * u
        bx = p3[0] + (p2[0]-p3[0]) * u
        by = p3[1] + (p2[1]-p3[1]) * u
        x = tx + (bx-tx) * v
        y = ty + (by-ty) * v
        min_d = 0.15 + 0.07 * rng.random()
        if any((x-qx)**2 + (y-qy)**2 < min_d**2 for qx, qy, *_ in accepted):
            continue
        # Outer-side bundles are a little smaller; inner-side bundles are larger.
        size = 0.17 + 0.10 * (0.30 + v) + rng.uniform(-0.016, 0.016)
        accepted.append((x, y, max(0.16, min(0.29, size)), rng.uniform(-15, 15)))
    return accepted


def add_grain_line(slide, x1, y1, x2, y2, color, width):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def render(config: dict, output: Path):
    prs = Presentation()
    prs.slide_width = Inches(float(config.get("slide_width", 13.333)))
    prs.slide_height = Inches(float(config.get("slide_height", 7.5)))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(config.get("background", "#FFFFFF"))

    top, end, side = projected_board_geometry(config)
    colors = config.get("colors", {})
    top_fill = colors.get("top", "#DCA234")
    side_fill = colors.get("side", "#C98B25")
    end_fill = colors.get("end", "#E5B13E")
    edge = colors.get("edge", "#9A6300")

    # Back-to-front layer order.
    add_polygon(slide, side, side_fill, edge, 0.8)
    add_polygon(slide, end, end_fill, edge, 0.8)
    add_polygon(slide, top, top_fill, edge, 0.8)

    rng = random.Random(int(config.get("seed", 20260903)))

    # Longitudinal grain on top.
    p0, p1, p2, p3 = top
    for i in range(int(config.get("top_grain_count", 30))):
        t = (i + 0.25 + rng.random()*0.5) / max(1, int(config.get("top_grain_count", 30)))
        sx = p0[0] + (p3[0]-p0[0])*t; sy = p0[1] + (p3[1]-p0[1])*t
        ex = p1[0] + (p2[0]-p1[0])*t; ey = p1[1] + (p2[1]-p1[1])*t
        add_grain_line(slide, sx, sy, ex, ey, colors.get("grain", "#B87917"), rng.uniform(0.20, 0.42))

    # Long-side grain.
    s0, s1, s2, s3 = side
    for i in range(int(config.get("side_grain_count", 10))):
        t = (i + 0.35 + rng.random()*0.3) / max(1, int(config.get("side_grain_count", 10)))
        sx = s0[0] + (s3[0]-s0[0])*t; sy = s0[1] + (s3[1]-s0[1])*t
        ex = s1[0] + (s2[0]-s1[0])*t; ey = s1[1] + (s2[1]-s1[1])*t
        add_grain_line(slide, sx, sy, ex, ey, colors.get("side_grain", "#A56F17"), rng.uniform(0.18, 0.34))

    # Fine editable speckles on transverse end face.
    a, b, c, d = end
    for _ in range(int(config.get("speckles", 85))):
        u, v = rng.random(), rng.random()
        tx = a[0] + (b[0]-a[0])*u; ty = a[1] + (b[1]-a[1])*u
        bx = d[0] + (c[0]-d[0])*u; by = d[1] + (c[1]-d[1])*u
        px = tx + (bx-tx)*v; py = ty + (by-ty)*v
        r = rng.uniform(0.005, 0.012)
        add_oval(slide, px-r, py-r, 2*r, 2*r, colors.get("speckle", "#9D6B18"))

    bundles = sample_nonuniform_bundles(rng, end, int(config.get("vascular_bundle_count", 18)))
    for cx, cy, size, rot in bundles:
        add_vascular_bundle(slide, cx, cy, size,
                            dark=colors.get("bundle_dark", "#6C3B05"),
                            gold=colors.get("bundle_gold", "#DDAE3E"),
                            rotation=rot)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    check = Presentation(output)
    shape_count = len(check.slides[0].shapes)
    if len(check.slides) != 1 or shape_count < 110 or len(bundles) < 12:
        raise RuntimeError("Bamboo board verification failed")
    return output, shape_count, len(bundles)


def main():
    ap = argparse.ArgumentParser(description="Render editable 5 x 2 x 0.5 bamboo board.")
    ap.add_argument("--config", type=Path)
    ap.add_argument("--output", required=True, type=Path)
    args = ap.parse_args()
    if args.config:
        cfg = json.loads(args.config.read_text(encoding="utf-8"))
    else:
        cfg = {
            "physical_cm": {"length": 5.0, "width": 2.0, "thickness": 0.5},
            "seed": 20260903,
            "vascular_bundle_count": 18
        }
    out, shapes, bundles = render(cfg, args.output.resolve())
    print(json.dumps({"output": str(out), "shapes": shapes, "vascular_bundles": bundles, "physical_cm": "5 x 2 x 0.5", "ratio": "10:4:1"}, ensure_ascii=False))


if __name__ == "__main__":
    main()
