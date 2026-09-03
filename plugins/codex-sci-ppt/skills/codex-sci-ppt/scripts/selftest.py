#!/usr/bin/env python3
"""End-to-end self-test for Codex Sci-PPT's local, no-API pipeline."""
from __future__ import annotations

import hashlib
import json
import subprocess
import sys
import tempfile
from pathlib import Path

import cv2
import numpy as np
from pptx import Presentation


def run(*args: object) -> str:
    completed = subprocess.run(
        [sys.executable, *map(str, args)],
        check=True,
        text=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
    )
    return completed.stdout.strip()


def sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def assert_pptx(path: Path, min_shapes: int = 1, expected_text: str | None = None) -> int:
    prs = Presentation(str(path))
    if not prs.slides:
        raise AssertionError(f"no slides in {path}")
    count = sum(len(slide.shapes) for slide in prs.slides)
    if count < min_shapes:
        raise AssertionError(f"expected >= {min_shapes} shapes, got {count} in {path}")
    if expected_text is not None:
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    texts.append(shape.text)
        if not any(expected_text in text for text in texts):
            raise AssertionError(f"editable text {expected_text!r} not found in {path}")
    return count


def main() -> int:
    scripts = Path(__file__).resolve().parent
    repo_root = Path(__file__).resolve().parents[5]
    summary: dict[str, object] = {"ok": False}

    with tempfile.TemporaryDirectory(prefix="codex-sci-ppt-selftest-") as td:
        work = Path(td)

        # 1) SVG validation/cache/cull/OOXML path, including transforms,
        # cubic geometry, duplicate paths, and editable text.
        svg = work / "core.svg"
        svg.write_text(
            '''<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 240 120">
  <g transform="translate(10 5) scale(2 1.5)">
    <rect id="rect-a" x="5" y="5" width="30" height="15" rx="3" fill="#D9EAD3" stroke="#38761D" stroke-width="2"/>
    <rect id="rect-dup" x="5" y="5" width="30" height="15" rx="3" fill="#D9EAD3" stroke="#38761D" stroke-width="2"/>
    <path id="curve" d="M 50 15 C 60 0 75 30 90 15 Z" fill="#F6B26B" stroke="#B45F06"/>
  </g>
  <text id="label" x="120" y="70" font-family="Arial" font-size="16" fill="#222222">Editable label</text>
</svg>\n''',
            encoding="utf-8",
        )
        run(scripts / "validate_vector_svg.py", "--svg", svg, "--strict-ids")
        cache_dir = work / "cache"
        run(scripts / "prepare_geometry_cache.py", "--input", svg, "--output-dir", cache_dir, "--job-id", "selftest")
        before = json.loads((cache_dir / "geometry-cache.json").read_text(encoding="utf-8"))
        if before["total_atoms"] != 4:
            raise AssertionError(f"expected 4 source atoms, got {before['total_atoms']}")
        first_anchor = before["atoms"][0]["subpaths"][0]["points"][0]["a"]
        if first_anchor == [8.0, 5.0] or first_anchor == [5.0, 5.0]:
            raise AssertionError("SVG transform was not applied to cached geometry")
        run(
            scripts / "cull_hidden_geometry.py",
            "--cache", cache_dir / "geometry-cache.json",
            "--state", cache_dir / "drawing-state.json",
        )
        after = json.loads((cache_dir / "geometry-cache.json").read_text(encoding="utf-8"))
        if after.get("culled_atom_count", 0) < 1:
            raise AssertionError("exact duplicate path was not culled")
        core_pptx = work / "core.pptx"
        run(scripts / "run_cell_ppt_ooxml.py", "--geometry-cache", cache_dir / "geometry-cache.json", "--output-pptx", core_pptx)
        core_shapes = assert_pptx(core_pptx, min_shapes=3, expected_text="Editable label")

        # 2) Scene renderer.
        scene_pptx = work / "scene.pptx"
        run(
            scripts / "render_scene.py",
            "--scene", repo_root / "examples" / "simple_scene.json",
            "--output", scene_pptx,
        )
        scene_shapes = assert_pptx(scene_pptx, min_shapes=3)

        # 3) Synthetic flat-color diagram used for local vectorizer and complete
        # image-pipeline tests.
        image = np.full((240, 360, 3), 255, dtype=np.uint8)
        cv2.rectangle(image, (30, 50), (150, 190), (230, 200, 120), thickness=-1)
        cv2.circle(image, (250, 120), 55, (90, 160, 235), thickness=-1)
        cv2.line(image, (150, 120), (195, 120), (60, 60, 60), thickness=8)
        cv2.putText(image, "TXT", (88, 35), cv2.FONT_HERSHEY_SIMPLEX, 0.8, (25, 25, 25), 2, cv2.LINE_AA)
        image_path = work / "synthetic.png"
        cv2.imwrite(str(image_path), image)

        # 3a) Vectorizer v2 should be deterministic, recover semantic primitives,
        # remain vector-only, and report a sensible palette reconstruction score.
        trace_one = work / "trace-one.svg"
        trace_two = work / "trace-two.svg"
        trace_args = [
            "--input-image", image_path,
            "--colors", "8",
            "--max-paths", "120",
            "--palette-merge-distance", "5",
        ]
        vector_one = json.loads(run(
            scripts / "vectorize_local.py", *trace_args,
            "--output-svg", trace_one,
        ).splitlines()[-1])
        vector_two = json.loads(run(
            scripts / "vectorize_local.py", *trace_args,
            "--output-svg", trace_two,
        ).splitlines()[-1])
        if sha256(trace_one) != sha256(trace_two):
            raise AssertionError("local vectorization is not deterministic")
        trace_payload = trace_one.read_text(encoding="utf-8")
        if "<image" in trace_payload:
            raise AssertionError("vectorizer emitted a raster image node")
        primitives = vector_one.get("primitives", {})
        if int(primitives.get("rect", 0)) < 1:
            raise AssertionError(f"expected rectangle primitive recovery, got {primitives}")
        if int(primitives.get("ellipse", 0)) < 1:
            raise AssertionError(f"expected ellipse primitive recovery, got {primitives}")
        if float(vector_one.get("palette_psnr_db", 0)) < 20:
            raise AssertionError(f"unexpectedly low palette PSNR: {vector_one.get('palette_psnr_db')}")
        if vector_one.get("palette_colors") != vector_two.get("palette_colors"):
            raise AssertionError("deterministic runs disagreed on palette size")

        # 3b) Full local image -> text cleanup -> SVG -> cache -> PPTX.
        manifest = work / "text-manifest.json"
        manifest.write_text(json.dumps({
            "schema_version": "1.0",
            "text_elements": [{
                "id": "live-text-1",
                "content": "TXT",
                "x": 0.245,
                "y": 0.145,
                "bbox": [0.225, 0.025, 0.19, 0.14],
                "coordinate_space": "normalized",
                "font_size": 18,
                "font_family": "Arial",
                "fill": "#191919",
                "paint_order": 999
            }]
        }), encoding="utf-8")
        image_pptx = work / "image.pptx"
        output = run(
            scripts / "run_from_image.py",
            "--input-image", image_path,
            "--text-manifest", manifest,
            "--output", image_pptx,
            "--colors", "8",
            "--max-paths", "120",
        )
        image_shapes = assert_pptx(image_pptx, min_shapes=2, expected_text="TXT")
        result = json.loads(output.splitlines()[-1])
        if int(result.get("text_regions_removed", 0)) != 1:
            raise AssertionError(f"expected one text region removed, got {result.get('text_regions_removed')}")
        if float(result.get("palette_psnr_db", 0)) < 20:
            raise AssertionError(f"pipeline did not propagate vectorizer quality metrics: {result}")
        master_svg = Path(result["master_svg"])
        master_payload = master_svg.read_text(encoding="utf-8")
        if "<image" in master_payload:
            raise AssertionError("local image pipeline left a raster <image> node in master SVG")
        if ">TXT<" not in master_payload:
            raise AssertionError("live editable text was not merged back into master SVG")

        # 3c) Transparency must stay transparent instead of becoming white art.
        rgba = np.zeros((120, 160, 4), dtype=np.uint8)
        cv2.circle(rgba, (80, 60), 35, (50, 120, 220, 255), thickness=-1)
        rgba_path = work / "transparent.png"
        cv2.imwrite(str(rgba_path), rgba)
        transparent_svg = work / "transparent.svg"
        transparent_result = json.loads(run(
            scripts / "vectorize_local.py",
            "--input-image", rgba_path,
            "--output-svg", transparent_svg,
            "--colors", "4",
        ).splitlines()[-1])
        transparent_payload = transparent_svg.read_text(encoding="utf-8")
        if "local_background" in transparent_payload:
            raise AssertionError("transparent input unexpectedly received an opaque background")
        if int(transparent_result.get("vector_elements", 0)) < 1:
            raise AssertionError("transparent input produced no vector artwork")

        # 4) Job allocator should produce stable sequential names.
        jobs = work / "jobs"
        first_job = subprocess.check_output(
            [sys.executable, str(scripts / "allocate_job_name.py"), "--root", str(jobs)], text=True
        ).strip()
        (jobs / first_job).mkdir(parents=True)
        second_job = subprocess.check_output(
            [sys.executable, str(scripts / "allocate_job_name.py"), "--root", str(jobs)], text=True
        ).strip()
        if (first_job, second_job) != ("codexscippt1", "codexscippt2"):
            raise AssertionError(f"unexpected job allocation: {first_job}, {second_job}")

        summary.update({
            "ok": True,
            "core_shapes": core_shapes,
            "scene_shapes": scene_shapes,
            "image_shapes": image_shapes,
            "source_atoms": before["total_atoms"],
            "retained_atoms": after["total_atoms"],
            "duplicates_removed": after.get("culled_atom_count", 0),
            "text_regions_removed": result.get("text_regions_removed", 0),
            "vectorizer_psnr_db": vector_one.get("palette_psnr_db"),
            "vectorizer_palette_colors": vector_one.get("palette_colors"),
            "vectorizer_primitives": primitives,
            "deterministic": True,
            "transparent_background_preserved": True,
            "raster_nodes": 0,
            "jobs": [first_job, second_job],
        })

    print(json.dumps(summary, ensure_ascii=False, separators=(",", ":")))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
