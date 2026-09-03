#!/usr/bin/env python3
"""Regression test for the local reference-image analyzer."""
from __future__ import annotations

import json
import tempfile
from pathlib import Path

import cv2
import numpy as np
from pptx import Presentation

from analyze_reference import analyze, build_scene_draft, _debug_overlay
from render_reference_scene import render


def synthetic_reference(path: Path) -> None:
    img = np.full((220, 322, 3), (248, 251, 252), np.uint8)
    teal = (155, 134, 20)      # BGR ~ #14869B
    magenta = (185, 38, 179)   # BGR
    yellow = (0, 170, 242)
    cyan = (198, 245, 248)
    dark = (35, 35, 35)

    cv2.rectangle(img, (21, 23), (265, 196), teal, 2)
    cv2.rectangle(img, (269, 35), (307, 189), teal, 2)
    cv2.rectangle(img, (75, 48), (139, 58), magenta, -1)
    cv2.rectangle(img, (164, 48), (231, 58), magenta, -1)
    cv2.ellipse(img, (103, 157), (21, 7), 0, 0, 360, teal, 1)
    cv2.ellipse(img, (157, 157), (21, 7), 0, 0, 360, teal, 1)
    for x in (92, 117, 180, 205):
        cv2.line(img, (x, 67), (x, 120), yellow, 2)
        cv2.fillConvexPoly(img, np.array([[x-4, 116], [x+4, 116], [x, 122]], np.int32), yellow)
    cv2.circle(img, (51, 102), 16, teal, 2)
    cv2.circle(img, (51, 102), 5, teal, -1)
    cv2.rectangle(img, (276, 61), (299, 79), (54, 45, 24), -1)
    cv2.rectangle(img, (276, 88), (299, 106), (54, 45, 24), -1)
    cv2.putText(img, "UVA-340 lamps", (113, 41), cv2.FONT_HERSHEY_SIMPLEX, 0.38, dark, 1, cv2.LINE_AA)
    cv2.putText(img, "Test samples", (128, 129), cv2.FONT_HERSHEY_SIMPLEX, 0.34, dark, 1, cv2.LINE_AA)
    cv2.putText(img, "IPBC", (93, 188), cv2.FONT_HERSHEY_SIMPLEX, 0.35, dark, 1, cv2.LINE_AA)
    if not cv2.imwrite(str(path), img):
        raise RuntimeError("failed to write synthetic reference")


def main() -> None:
    with tempfile.TemporaryDirectory(prefix="codex-sci-ppt-reference-") as td:
        root = Path(td)
        source = root / "reference.png"
        synthetic_reference(source)

        a1, rgb1 = analyze(source, threshold=14.0, palette_colors=8)
        a2, _ = analyze(source, threshold=14.0, palette_colors=8)
        if json.dumps(a1, sort_keys=True) != json.dumps(a2, sort_keys=True):
            raise AssertionError("reference analysis must be deterministic")
        if a1["image"]["width"] != 322 or a1["image"]["height"] != 220:
            raise AssertionError("image dimensions were not preserved")
        if len(a1["palette"]) < 3:
            raise AssertionError("expected at least three foreground palette colors")
        if not a1["frames"]:
            raise AssertionError("expected at least one large frame candidate")
        if a1["summary"]["high_confidence_primitive_count"] < 4:
            raise AssertionError("expected several high-confidence editable primitives")
        if a1["summary"]["text_like_region_count"] < 1:
            raise AssertionError("expected at least one text-like region")

        scene = build_scene_draft(a1, 10.0)
        if len(scene["objects"]) < 4:
            raise AssertionError("scene draft did not contain enough editable primitives")
        pptx = root / "draft.pptx"
        render(scene, pptx)
        check = Presentation(pptx)
        if not check.slides or len(check.slides[0].shapes) < 4:
            raise AssertionError("draft scene PPTX did not reopen with editable shapes")

        overlay = _debug_overlay(rgb1, a1)
        overlay_path = root / "overlay.png"
        if not cv2.imwrite(str(overlay_path), overlay) or not overlay_path.exists():
            raise AssertionError("debug overlay was not written")

        print(json.dumps({
            "ok": True,
            "background": a1["background"],
            "palette_colors": len(a1["palette"]),
            "components": a1["summary"]["component_count"],
            "high_confidence_primitives": a1["summary"]["high_confidence_primitive_count"],
            "frames": a1["summary"]["frame_candidate_count"],
            "text_like_regions": a1["summary"]["text_like_region_count"],
            "draft_shapes": len(check.slides[0].shapes),
        }, separators=(",", ":")))


if __name__ == "__main__":
    main()
