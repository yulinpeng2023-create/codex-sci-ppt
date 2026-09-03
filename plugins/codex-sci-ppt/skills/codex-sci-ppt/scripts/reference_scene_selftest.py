#!/usr/bin/env python3
from __future__ import annotations

import json
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation


def main() -> int:
    scripts = Path(__file__).resolve().parent
    repo_root = Path(__file__).resolve().parents[5]
    scene = repo_root / 'examples' / 'uva_chamber_scene.json'
    with tempfile.TemporaryDirectory(prefix='codex-sci-ppt-reference-scene-') as td:
        out = Path(td) / 'reference-scene.pptx'
        subprocess.run([sys.executable, str(scripts / 'render_reference_scene.py'),
                        '--scene', str(scene), '--output', str(out)], check=True)
        prs = Presentation(out)
        shapes = list(prs.slides[0].shapes)
        if len(shapes) < 30:
            raise AssertionError(f'expected >=30 editable shapes, got {len(shapes)}')
        texts = [s.text for s in shapes if getattr(s, 'has_text_frame', False)]
        for expected in ('UVA-340 lamps', 'Exhaust fan', 'Test samples', 'IPBC', 'HDI', 'Coated bamboo', '240 h', 'RH', '30 cm'):
            if not any(expected in t for t in texts):
                raise AssertionError(f'missing editable text: {expected}')
        print(json.dumps({'ok': True, 'shapes': len(shapes), 'texts_checked': 9}, separators=(',', ':')))
    return 0


if __name__ == '__main__':
    raise SystemExit(main())
