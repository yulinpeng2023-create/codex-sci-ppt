#!/usr/bin/env python3
"""Render a reusable bamboo cross-section schematic template.

The template is based on the user-supplied PowerPoint visual language: a warm
golden parenchyma matrix with repeated stylized vascular bundles. The matrix
texture is generated locally and deterministically; vascular bundles are native
editable PowerPoint shapes.
"""
from __future__ import annotations

import argparse
import json
import random
import tempfile
from pathlib import Path

import numpy as np
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#").upper())


def add_shape(slide, preset, x, y, w, h, fill, line=None, line_width=0.5, rotation=0):
    shape = slide.shapes.add_shape(preset, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.rotation = float(rotation)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(float(line_width))
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
    return shape


def make_matrix_texture(path: Path, width=1500, height=500, seed=20260903, base="#E5B13E"):
    rng = np.random.default_rng(seed)
    base_rgb = tuple(int(base[i:i + 2], 16) for i in (1, 3, 5))
    arr = np.zeros((height, width, 3), dtype=np.float32)
    arr[:] = base_rgb

    low = rng.normal(0, 1, (max(2, height // 32), max(2, width // 32))).astype(np.float32)
    low_u8 = np.uint8(np.clip((low - low.min()) / (low.max() - low.min() + 1e-6) * 255, 0, 255))
    low_img = Image.fromarray(low_u8).resize((width, height), Image.Resampling.BICUBIC)
    low_img = low_img.filter(ImageFilter.GaussianBlur(radius=4))
    arr += (np.asarray(low_img, dtype=np.float32) - 127.5)[..., None] * 0.10
    arr += rng.normal(0, 11, (height, width, 1)).astype(np.float32)
    arr = np.uint8(np.clip(arr, 0, 255))

    image = Image.fromarray(arr, "RGB")
    draw = ImageDraw.Draw(image, "RGBA")
    n = int(width * height * 0.012)
    xs = rng.integers(0, width, n)
    ys = rng.integers(0, height, n)
    rs = rng.integers(1, 3, n)
    for x, y, radius in zip(xs, ys, rs):
        draw.ellipse(
            (int(x - radius), int(y - radius), int(x + radius), int(y + radius)),
            fill=(80, 52, 14, int(rng.integers(18, 55))),
        )
    image = image.filter(ImageFilter.GaussianBlur(radius=0.35))
    image.save(path, format="PNG", optimize=True)


def add_vascular_bundle(slide, cx, cy, size=0.30, rotation=0, dark="#5A2F00", gold="#E3B247"):
    """Draw one stylized bamboo vascular bundle as editable PowerPoint shapes."""
    add_shape(slide, MSO_SHAPE.OVAL, cx - size * 0.18, cy - size * 0.46,
              size * 0.36, size * 0.36, dark, None, rotation=rotation)

    for sign in (-1, 1):
        lx = cx + sign * size * 0.30
        ly = cy - size * 0.02
        add_shape(slide, MSO_SHAPE.OVAL, lx - size * 0.20, ly - size * 0.20,
                  size * 0.40, size * 0.40, dark, None, rotation=rotation)
        add_shape(slide, MSO_SHAPE.OVAL, lx - size * 0.12, ly - size * 0.12,
                  size * 0.24, size * 0.24, gold, None, rotation=rotation)

    add_shape(slide, MSO_SHAPE.TEAR, cx - size * 0.18, cy + size * 0.10,
              size * 0.36, size * 0.44, dark, None, rotation=180 + rotation)
    add_shape(slide, MSO_SHAPE.OVAL, cx - size * 0.18, cy - size * 0.13,
              size * 0.36, size * 0.36, gold, None, rotation=rotation)

    for dx in (-0.055, 0.055):
        add_shape(slide, MSO_SHAPE.OVAL,
                  cx + size * dx - size * 0.018, cy - size * 0.02,
                  size * 0.036, size * 0.036, dark, None, rotation=rotation)


def render(config: dict, output: Path) -> tuple[Path, int]:
    prs = Presentation()
    prs.slide_width = Inches(float(config["slide"]["width"]))
    prs.slide_height = Inches(float(config["slide"]["height"]))
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb(config.get("slide_background", "#FFFFFF"))

    region = config["region"]
    x, y, w, h = map(float, (region["x"], region["y"], region["w"], region["h"]))

    with tempfile.TemporaryDirectory() as temp_dir:
        texture = Path(temp_dir) / "bamboo-matrix.png"
        make_matrix_texture(
            texture,
            width=int(config.get("texture_width", 1500)),
            height=int(config.get("texture_height", 500)),
            seed=int(config.get("seed", 20260903)),
            base=config.get("matrix_color", "#E5B13E"),
        )
        slide.shapes.add_picture(str(texture), Inches(x), Inches(y), Inches(w), Inches(h))

    rng = random.Random(int(config.get("seed", 20260903)))
    count = int(config.get("vascular_bundle_count", 55))
    margin = float(config.get("bundle_margin", 0.18))
    min_distance = float(config.get("min_distance", 0.58))
    points = []
    attempts = 0

    while len(points) < count and attempts < count * 500:
        attempts += 1
        px = x + margin + rng.random() * (w - 2 * margin)
        py = y + margin + rng.random() * (h - 2 * margin)
        relative_y = (py - y) / h
        outer_bias = float(config.get("outer_density_bias", 0.12))
        if outer_bias > 0 and rng.random() > 1.0 - outer_bias * (1 - relative_y):
            continue
        if any((px - qx) ** 2 + (py - qy) ** 2 < min_distance ** 2 for qx, qy, _ in points):
            continue
        size = rng.uniform(
            float(config.get("bundle_size_min", 0.26)),
            float(config.get("bundle_size_max", 0.34)),
        )
        points.append((px, py, size))

    while len(points) < count:
        px = x + margin + rng.random() * (w - 2 * margin)
        py = y + margin + rng.random() * (h - 2 * margin)
        size = rng.uniform(0.26, 0.34)
        points.append((px, py, size))

    for px, py, size in points:
        add_vascular_bundle(
            slide,
            px,
            py,
            size=size,
            rotation=rng.uniform(-18, 18),
            dark=config.get("bundle_dark", "#5A2F00"),
            gold=config.get("bundle_gold", "#E3B247"),
        )

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)

    check = Presentation(output)
    shape_count = len(check.slides[0].shapes)
    if len(check.slides) != 1 or shape_count < count * 7:
        raise RuntimeError("Bamboo cross-section template verification failed")
    return output, shape_count


def main():
    parser = argparse.ArgumentParser(description="Render the reusable bamboo cross-section template.")
    parser.add_argument("--config", type=Path)
    parser.add_argument("--output", required=True, type=Path)
    args = parser.parse_args()

    if args.config:
        config = json.loads(args.config.read_text(encoding="utf-8"))
    else:
        config = {
            "slide": {"width": 13.333, "height": 7.5},
            "slide_background": "#FFFFFF",
            "region": {"x": 0.55, "y": 1.65, "w": 12.25, "h": 3.95},
            "matrix_color": "#E5B13E",
            "vascular_bundle_count": 55,
            "seed": 20260903,
            "bundle_dark": "#5A2F00",
            "bundle_gold": "#E3B247",
            "min_distance": 0.58,
            "bundle_size_min": 0.26,
            "bundle_size_max": 0.34,
            "outer_density_bias": 0.12,
        }

    output, shape_count = render(config, args.output.resolve())
    print(json.dumps({"output": str(output), "shapes": shape_count}, ensure_ascii=False))


if __name__ == "__main__":
    main()
