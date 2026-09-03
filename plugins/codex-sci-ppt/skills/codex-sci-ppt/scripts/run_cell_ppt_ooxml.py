#!/usr/bin/env python3
"""Cross-platform editable OOXML renderer for Codex Sci-PPT geometry caches.

Adapted from the MIT-licensed yrui-cmd/cell-ppt renderer. It appends native
custom-geometry shapes and text boxes to a PPTX without raster fallback.
See THIRD_PARTY_NOTICES.md.
"""

from __future__ import annotations

import argparse
import json
import os
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Emu, Pt

EMU_PER_PT = 12700.0
PATH_EXTENT = 100000


def esc(value: str) -> str:
    return (
        value.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
        .replace("'", "&apos;")
    )


def same(a: list[float], b: list[float]) -> bool:
    return abs(a[0] - b[0]) < 1e-6 and abs(a[1] - b[1]) < 1e-6


def rgb_hex(value: list[int] | None, fallback: str = "000000") -> str:
    if not value:
        return fallback
    return "".join(f"{max(0, min(255, int(v))):02X}" for v in value[:3])


def alpha_xml(opacity: float) -> str:
    alpha = max(0, min(100000, round(float(opacity) * 1000)))
    return f'<a:alpha val="{alpha}"/>' if alpha < 100000 else ""


def map_point(point, scale, offset_x, offset_y, view_x, view_y):
    return (
        offset_x + (float(point[0]) - view_x) * scale,
        offset_y + (float(point[1]) - view_y) * scale,
    )


def add_freeform(slide, subpath, paint, name, shape_id, transform):
    scale, offset_x, offset_y, view_x, view_y = transform
    points = subpath["points"]
    mapped = []
    for point in points:
        mapped.append(
            {
                "a": map_point(point["a"], scale, offset_x, offset_y, view_x, view_y),
                "l": map_point(point["l"], scale, offset_x, offset_y, view_x, view_y),
                "r": map_point(point["r"], scale, offset_x, offset_y, view_x, view_y),
            }
        )
    all_xy = [coord for point in mapped for key in ("a", "l", "r") for coord in [point[key]]]
    min_x = min(p[0] for p in all_xy)
    min_y = min(p[1] for p in all_xy)
    max_x = max(p[0] for p in all_xy)
    max_y = max(p[1] for p in all_xy)
    width = max(max_x - min_x, 1.0 / EMU_PER_PT)
    height = max(max_y - min_y, 1.0 / EMU_PER_PT)

    def local(point):
        x = round((point[0] - min_x) / width * PATH_EXTENT)
        y = round((point[1] - min_y) / height * PATH_EXTENT)
        return max(-2147483647, min(2147483647, x)), max(-2147483647, min(2147483647, y))

    first_x, first_y = local(mapped[0]["a"])
    commands = [f'<a:moveTo><a:pt x="{first_x}" y="{first_y}"/></a:moveTo>']
    for idx in range(1, len(mapped)):
        previous = mapped[idx - 1]
        current = mapped[idx]
        end_x, end_y = local(current["a"])
        if same(previous["r"], previous["a"]) and same(current["l"], current["a"]):
            commands.append(f'<a:lnTo><a:pt x="{end_x}" y="{end_y}"/></a:lnTo>')
        else:
            c1x, c1y = local(previous["r"])
            c2x, c2y = local(current["l"])
            commands.append(
                f'<a:cubicBezTo><a:pt x="{c1x}" y="{c1y}"/>'
                f'<a:pt x="{c2x}" y="{c2y}"/><a:pt x="{end_x}" y="{end_y}"/></a:cubicBezTo>'
            )
    if subpath.get("closed"):
        previous = mapped[-1]
        current = mapped[0]
        if not (same(previous["r"], previous["a"]) and same(current["l"], current["a"])):
            c1x, c1y = local(previous["r"])
            c2x, c2y = local(current["l"])
            commands.append(
                f'<a:cubicBezTo><a:pt x="{c1x}" y="{c1y}"/>'
                f'<a:pt x="{c2x}" y="{c2y}"/><a:pt x="{first_x}" y="{first_y}"/></a:cubicBezTo>'
            )
        commands.append("<a:close/>")

    fill = "<a:noFill/>"
    if paint.get("filled") and subpath.get("closed"):
        fill = (
            f'<a:solidFill><a:srgbClr val="{rgb_hex(paint.get("fillColor"))}">'
            f'{alpha_xml(paint.get("opacity", 100))}</a:srgbClr></a:solidFill>'
        )
    line = '<a:ln><a:noFill/></a:ln>'
    if paint.get("stroked"):
        weight = max(3175, round(float(paint.get("strokeWidth", 1)) * scale * EMU_PER_PT))
        line = (
            f'<a:ln w="{weight}"><a:solidFill><a:srgbClr val="{rgb_hex(paint.get("strokeColor"))}">'
            f'{alpha_xml(paint.get("opacity", 100))}</a:srgbClr></a:solidFill></a:ln>'
        )
    off_x, off_y = round(min_x * EMU_PER_PT), round(min_y * EMU_PER_PT)
    ext_x, ext_y = max(1, round(width * EMU_PER_PT)), max(1, round(height * EMU_PER_PT))
    xml = f'''<p:sp {nsdecls("a", "p")}>
      <p:nvSpPr><p:cNvPr id="{shape_id}" name="{esc(name)}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="{off_x}" y="{off_y}"/><a:ext cx="{ext_x}" cy="{ext_y}"/></a:xfrm>
        <a:custGeom><a:avLst/><a:gdLst/><a:ahLst/><a:cxnLst/><a:rect l="l" t="t" r="r" b="b"/>
          <a:pathLst><a:path w="{PATH_EXTENT}" h="{PATH_EXTENT}">{''.join(commands)}</a:path></a:pathLst>
        </a:custGeom>{fill}{line}
      </p:spPr>
      <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
    </p:sp>'''
    slide.shapes._spTree.insert_element_before(parse_xml(xml), "p:extLst")


def add_text(slide, atom, transform):
    scale, offset_x, offset_y, view_x, view_y = transform
    text = atom["text"]
    x, y = map_point(text["position"], scale, offset_x, offset_y, view_x, view_y)
    size = max(4.0, float(text["fontSize"]) * scale)
    width = max(size * 2.0, size * len(str(text["contents"])) * 0.7)
    height = size * 1.5
    anchor = str(text.get("textAnchor") or "start").lower()
    if anchor == "middle":
        x -= width / 2.0
    elif anchor == "end":
        x -= width
    shape = slide.shapes.add_textbox(
        Emu(round(x * EMU_PER_PT)),
        Emu(round((y - size * 1.05) * EMU_PER_PT)),
        Emu(round(width * EMU_PER_PT)),
        Emu(round(height * EMU_PER_PT)),
    )
    shape.name = atom["objectName"]
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = str(text["contents"])
    font = run.font
    font.name = str(text.get("fontFamily") or "Arial")
    font.size = Pt(size)
    weight = str(text.get("fontWeight", "")).strip().lower()
    font.bold = weight == "bold" or (weight.isdigit() and int(weight) >= 600)
    font.italic = str(text.get("fontStyle", "")).lower() == "italic"
    font.color.rgb = RGBColor(*[int(v) for v in text.get("fillColor", [0, 0, 0])])
    shape.rotation = float(text.get("rotationDegrees", 0))


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--geometry-cache", required=True, type=Path)
    parser.add_argument("--output-pptx", required=True, type=Path)
    parser.add_argument("--input-pptx", type=Path)
    parser.add_argument("--slide-index", type=int, default=0, help="1-based existing slide; 0 appends a blank slide")
    args = parser.parse_args()
    cache = json.loads(args.geometry_cache.read_text(encoding="utf-8"))
    if int(cache.get("schema_version", 0)) != 3:
        raise SystemExit("Unsupported geometry cache schema")
    prs = Presentation(str(args.input_pptx)) if args.input_pptx else Presentation()
    if args.slide_index:
        if not 1 <= args.slide_index <= len(prs.slides):
            raise SystemExit("slide-index is out of range")
        slide = prs.slides[args.slide_index - 1]
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width / EMU_PER_PT
    slide_height = prs.slide_height / EMU_PER_PT
    view_x, view_y, view_width, view_height = map(float, cache["view_box"])
    margin = 18.0
    scale = min((slide_width - 2 * margin) / view_width, (slide_height - 2 * margin) / view_height)
    offset_x = (slide_width - view_width * scale) / 2.0
    offset_y = (slide_height - view_height * scale) / 2.0
    transform = (scale, offset_x, offset_y, view_x, view_y)
    next_id = max((shape.shape_id for shape in slide.shapes), default=1) + 1
    created = 0
    for batch in cache["batches"]:
        for atom_index in batch["atom_indices"]:
            atom = cache["atoms"][int(atom_index)]
            if atom["kind"] == "text":
                add_text(slide, atom, transform)
                created += 1
                next_id += 1
            elif atom["kind"] == "path":
                for part_index, subpath in enumerate(atom.get("subpaths", [])):
                    if len(subpath.get("points", [])) < 2:
                        continue
                    paints = atom.get("paintParts") or [{}]
                    paint = paints[min(part_index, len(paints) - 1)]
                    add_freeform(
                        slide,
                        subpath,
                        paint,
                        f'{atom["objectName"]}_PART_{part_index:03d}',
                        next_id,
                        transform,
                    )
                    next_id += 1
                    created += 1
    args.output_pptx.parent.mkdir(parents=True, exist_ok=True)
    handle, temporary_name = tempfile.mkstemp(
        prefix=f".{args.output_pptx.stem}-",
        suffix=".pptx",
        dir=args.output_pptx.parent,
    )
    os.close(handle)
    temporary = Path(temporary_name)
    try:
        prs.save(str(temporary))
        Presentation(str(temporary))
        os.replace(temporary, args.output_pptx)
    finally:
        if temporary.exists():
            temporary.unlink()
    reopened = Presentation(str(args.output_pptx))
    print(json.dumps({
        "ok": True,
        "backend": "editable-ooxml",
        "output_pptx": str(args.output_pptx.resolve()),
        "slide_index": args.slide_index or len(reopened.slides),
        "native_object_count": created,
        "slide_count": len(reopened.slides),
    }, ensure_ascii=False, separators=(",", ":")))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
