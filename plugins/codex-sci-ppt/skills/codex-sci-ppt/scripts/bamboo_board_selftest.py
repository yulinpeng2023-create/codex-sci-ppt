#!/usr/bin/env python3
"""Regression checks for the 5 x 2 x 0.5 cm editable bamboo board."""
from __future__ import annotations

import json
import math
import random
import tempfile
from pathlib import Path

from pptx import Presentation

from render_bamboo_board import board_geometry, render, sample_bundles


def distance(a, b):
    return math.hypot(b[0] - a[0], b[1] - a[1])


def main():
    root = Path(__file__).resolve().parents[1]
    config_path = root / "templates" / "bamboo_board_5x2x05.json"
    cfg = json.loads(config_path.read_text(encoding="utf-8"))

    assert cfg["physical_cm"] == {"length": 5.0, "width": 2.0, "thickness": 0.5}

    g = board_geometry(cfg)
    front = g["front"]
    # The transverse end is intentionally presented face-on so the 2:0.5
    # physical aspect is visually preserved as 4:1.
    assert abs(distance(front[0], front[1]) / distance(front[0], front[3]) - 4.0) < 0.02
    assert len(g["top"]) == 4 and len(g["side"]) == 4

    bundles = sample_bundles(random.Random(cfg["seed"]), front, cfg["vascular_bundle_count"])
    assert len(bundles) == cfg["vascular_bundle_count"]
    xs = [p[0] for p, _, _, _ in bundles]
    # Guard against the old evenly spaced/grid-like regression.
    gaps = sorted(round(xs[i + 1] - xs[i], 3) for i in range(len(xs) - 1)) if len(xs) > 1 else []
    assert len(set(gaps)) >= 3

    with tempfile.TemporaryDirectory() as td:
        out = Path(td) / "bamboo-board.pptx"
        shape_count, bundle_count = render(cfg, out)
        check = Presentation(out)
        assert len(check.slides) == 1
        assert bundle_count == cfg["vascular_bundle_count"]
        assert shape_count >= 100
        assert len(check.slides[0].shapes) == shape_count

    print("bamboo board self-test: OK")


if __name__ == "__main__":
    main()
