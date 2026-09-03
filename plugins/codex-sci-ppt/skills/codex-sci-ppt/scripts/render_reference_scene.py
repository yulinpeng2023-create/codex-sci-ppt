#!/usr/bin/env python3
"""High-fidelity editable scene renderer for flat scientific reference diagrams.

This is an additive scene-mode renderer. It does not change the Cell-PPT-
compatible raster/SVG/cache/OOXML reconstruction pipeline.
"""
from __future__ import annotations

import argparse
import json
import math
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip('#').upper())


def strip_theme_style(shape) -> None:
    try:
        style = shape._element.find(qn('p:style'))
        if style is not None:
            shape._element.remove(style)
        shape.shadow.inherit = False
    except Exception:
        pass


def style_shape(shape, fill='#FFFFFF', line='#333333', line_width=1.0, rotation=0.0):
    strip_theme_style(shape)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(float(line_width))
    shape.rotation = float(rotation)
    return shape


def line_end(shape, tag: str, kind: str | None, width='sm', length='sm') -> None:
    if not kind or kind == 'none':
        return
    ln = shape.line._get_or_add_ln()
    node = OxmlElement(f'a:{tag}')
    node.set('type', kind)
    node.set('w', width)
    node.set('len', length)
    ln.append(node)


def add_text(slide, o):
    sh = slide.shapes.add_textbox(Inches(o['x']), Inches(o['y']), Inches(o['w']), Inches(o['h']))
    strip_theme_style(sh)
    sh.rotation = float(o.get('rotation', 0))
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = bool(o.get('word_wrap', True))
    m = Inches(float(o.get('margin', 0)))
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = m
    tf.vertical_anchor = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE,
                          'bottom': MSO_ANCHOR.BOTTOM}.get(o.get('valign', 'middle'), MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.text = str(o.get('text', ''))
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER,
                   'right': PP_ALIGN.RIGHT}.get(o.get('align', 'left'), PP_ALIGN.LEFT)
    p.space_before = p.space_after = Pt(0)
    if p.runs:
        r = p.runs[0]
        r.font.name = o.get('font', 'Arial')
        r.font.size = Pt(float(o.get('font_size', 10)))
        r.font.bold = bool(o.get('bold', False))
        r.font.italic = bool(o.get('italic', False))
        r.font.color.rgb = rgb(o.get('color', '#222222'))
    return sh


def add_shape(slide, o, preset):
    sh = slide.shapes.add_shape(preset, Inches(o['x']), Inches(o['y']), Inches(o['w']), Inches(o['h']))
    return style_shape(sh, o.get('fill', '#FFFFFF'), o.get('line', '#333333'),
                       o.get('line_width', 1.0), o.get('rotation', 0))


def add_polygon(slide, points, fill, line, line_width=0.7):
    b = slide.shapes.build_freeform(Inches(points[0][0]), Inches(points[0][1]))
    b.add_line_segments([(Inches(x), Inches(y)) for x, y in points[1:]], close=True)
    sh = b.convert_to_shape()
    return style_shape(sh, fill, line, line_width)


def add_connector(slide, o, arrow=False):
    sh = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(o['x1']), Inches(o['y1']),
                                    Inches(o['x2']), Inches(o['y2']))
    strip_theme_style(sh)
    sh.line.color.rgb = rgb(o.get('line', '#333333'))
    sh.line.width = Pt(float(o.get('line_width', 1.0)))
    if arrow:
        line_end(sh, 'headEnd', o.get('start_arrowhead'), o.get('arrow_width', 'sm'), o.get('arrow_length', 'sm'))
        line_end(sh, 'tailEnd', o.get('end_arrowhead', o.get('arrowhead', 'triangle')),
                 o.get('arrow_width', 'sm'), o.get('arrow_length', 'sm'))
    return sh


def label(slide, text, x, y, w, h, **kw):
    o = {'x': x, 'y': y, 'w': w, 'h': h, 'text': text, 'align': 'center', 'valign': 'middle',
         'font': kw.get('font', 'Times New Roman'), 'font_size': kw.get('font_size', 9),
         'bold': kw.get('bold', False), 'color': kw.get('color', '#333333'),
         'rotation': kw.get('rotation', 0), 'margin': 0}
    return add_text(slide, o)


def uv_lamp(slide, o):
    add_shape(slide, o, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_shape(slide, {'x': o['x'] + o['w']*0.05, 'y': o['y'] + o['h']*0.14,
                      'w': o['w']*0.90, 'h': max(o['h']*0.16, 0.015),
                      'fill': o.get('highlight', '#DA70DD'), 'line': None}, MSO_SHAPE.ROUNDED_RECTANGLE)


def fan(slide, o):
    x, y, w, h = o['x'], o['y'], o['w'], o['h']
    add_shape(slide, {'x': x, 'y': y, 'w': w, 'h': h, 'fill': o.get('fill', '#F8FCFD'),
                      'line': o.get('line', '#2D93A6'), 'line_width': o.get('line_width', 0.8)}, MSO_SHAPE.OVAL)
    cx, cy = x+w/2, y+h/2
    n = int(o.get('blades', 3))
    bw, bh = w*0.30, h*0.42
    for i in range(n):
        angle = -90 + i*360/n
        rad = math.radians(angle)
        px = cx + math.cos(rad)*w*0.17 - bw/2
        py = cy + math.sin(rad)*h*0.17 - bh/2
        add_shape(slide, {'x': px, 'y': py, 'w': bw, 'h': bh, 'fill': o.get('blade_fill', '#78B9C8'),
                          'line': None, 'rotation': angle+90}, MSO_SHAPE.TEAR)
    add_shape(slide, {'x': cx-w*0.075, 'y': cy-h*0.075, 'w': w*0.15, 'h': h*0.15,
                      'fill': '#FFFFFF', 'line': o.get('line', '#2D93A6'), 'line_width': 0.4}, MSO_SHAPE.OVAL)
    if o.get('label'):
        label(slide, o['label'], x-w*0.35, y+h+0.03, w*1.7, 0.28,
              font_size=o.get('font_size', 8), color=o.get('text_color', '#666666'))


def petri_dish(slide, o):
    x, y, w, h = o['x'], o['y'], o['w'], o['h']
    line = o.get('line', '#20889A')
    lw = o.get('line_width', 0.6)
    add_shape(slide, {'x': x, 'y': y+h*0.24, 'w': w, 'h': h*0.48,
                      'fill': o.get('fill', '#E8F5F6'), 'line': line, 'line_width': lw}, MSO_SHAPE.OVAL)
    add_shape(slide, {'x': x, 'y': y, 'w': w, 'h': h*0.48,
                      'fill': '#FFFFFF', 'line': line, 'line_width': lw}, MSO_SHAPE.OVAL)
    add_connector(slide, {'x1': x, 'y1': y+h*0.24, 'x2': x, 'y2': y+h*0.49, 'line': line, 'line_width': lw})
    add_connector(slide, {'x1': x+w, 'y1': y+h*0.24, 'x2': x+w, 'y2': y+h*0.49, 'line': line, 'line_width': lw})
    if o.get('inner_fill'):
        add_shape(slide, {'x': x+w*0.10, 'y': y+h*0.07, 'w': w*0.80, 'h': h*0.30,
                          'fill': o['inner_fill'], 'line': None}, MSO_SHAPE.OVAL)
    if o.get('label'):
        label(slide, o['label'], x-w*0.18, y+h*0.80, w*1.36, 0.28,
              font_size=o.get('font_size', 8.6), bold=o.get('bold', True))


def material_block(slide, o):
    x, y, w, h = o['x'], o['y'], o['w'], o['h']
    depth = o.get('depth', min(w*0.22, h*0.55))
    skew = o.get('skew', depth*0.62)
    top = [[x, y+skew], [x+depth, y], [x+w, y], [x+w-depth, y+skew]]
    front = [[x, y+skew], [x+w-depth, y+skew], [x+w-depth, y+h], [x, y+h]]
    side = [[x+w-depth, y+skew], [x+w, y], [x+w, y+h-skew], [x+w-depth, y+h]]
    add_polygon(slide, front, o.get('front_fill', '#D0A756'), o.get('line', '#7B6937'), 0.6)
    add_polygon(slide, side, o.get('side_fill', '#9D6A2F'), o.get('line', '#7B6937'), 0.6)
    add_polygon(slide, top, o.get('top_fill', '#D8C55F'), o.get('line', '#7B6937'), 0.6)
    for i in range(1, int(o.get('grain_lines', 4))+1):
        t = i/(int(o.get('grain_lines', 4))+1)
        add_connector(slide, {'x1': x+depth*t, 'y1': y+skew*(1-t),
                              'x2': x+w-depth*(1-t), 'y2': y+skew*t*0.08,
                              'line': o.get('grain_line', '#A48B3C'), 'line_width': 0.28})
    if o.get('label'):
        label(slide, o['label'], x-0.15, y+h+0.07, w+0.30, 0.30,
              font_size=o.get('font_size', 8.3), bold=o.get('bold', True))


def dimension_arrow(slide, o):
    add_connector(slide, {'x1': o['x1'], 'y1': o['y1'], 'x2': o['x2'], 'y2': o['y2'],
                          'line': o.get('line', '#888888'), 'line_width': o.get('line_width', 0.55),
                          'start_arrowhead': o.get('start_arrowhead', 'triangle'),
                          'end_arrowhead': o.get('end_arrowhead', 'triangle'),
                          'arrow_width': 'sm', 'arrow_length': 'sm'}, arrow=True)
    if o.get('label'):
        mx, my = (o['x1']+o['x2'])/2, (o['y1']+o['y2'])/2
        label(slide, o['label'], mx-0.18, my-0.55, 0.36, 1.10,
              font_size=o.get('font_size', 8), color=o.get('text_color', '#777777'), rotation=270)


def control_panel(slide, o):
    x, y, w, h = o['x'], o['y'], o['w'], o['h']
    edge = o.get('line', '#14869B')
    add_shape(slide, {'x': x, 'y': y, 'w': w, 'h': h, 'fill': o.get('fill', '#EAF0F2'),
                      'line': edge, 'line_width': o.get('line_width', 0.8)}, MSO_SHAPE.RECTANGLE)
    dw, dh = w*0.68, h*0.11
    dx = x+(w-dw)/2
    for i, text in enumerate(o.get('displays', ['240 h', 'RH'])[:2]):
        yy = y+h*0.17+i*h*0.18
        add_shape(slide, {'x': dx, 'y': yy, 'w': dw, 'h': dh, 'fill': '#17333B', 'line': None}, MSO_SHAPE.RECTANGLE)
        label(slide, text, dx, yy, dw, dh, font='Arial', font_size=o.get('display_font_size', 8), bold=True, color='#FFFFFF')
    d = w*0.075
    for i, color in enumerate(o.get('indicators', ['#F1AA00', '#4C9CB5', '#D15A54'])[:3]):
        add_shape(slide, {'x': x+w*0.26+i*w*0.20, 'y': y+h*0.58, 'w': d, 'h': d,
                          'fill': color, 'line': '#7E7E7E', 'line_width': 0.2}, MSO_SHAPE.OVAL)
    cd = w*0.30
    cx, cy = x+(w-cd)/2, y+h*0.76
    add_shape(slide, {'x': cx, 'y': cy, 'w': cd, 'h': cd, 'fill': '#FFFFFF',
                      'line': '#4E5C63', 'line_width': 0.6}, MSO_SHAPE.OVAL)
    ccx, ccy = cx+cd/2, cy+cd/2
    add_connector(slide, {'x1': ccx, 'y1': ccy, 'x2': ccx+cd*0.18, 'y2': ccy-cd*0.12, 'line': '#4E5C63', 'line_width': 0.5})
    add_connector(slide, {'x1': ccx, 'y1': ccy, 'x2': ccx, 'y2': ccy-cd*0.23, 'line': '#4E5C63', 'line_width': 0.5})


def render_object(slide, o):
    t = o['type']
    if t == 'text': return add_text(slide, o)
    if t == 'rect': return add_shape(slide, o, MSO_SHAPE.RECTANGLE)
    if t == 'round_rect': return add_shape(slide, o, MSO_SHAPE.ROUNDED_RECTANGLE)
    if t == 'ellipse': return add_shape(slide, o, MSO_SHAPE.OVAL)
    if t == 'line': return add_connector(slide, o)
    if t == 'arrow': return add_connector(slide, o, arrow=True)
    if t == 'uv_lamp': return uv_lamp(slide, o)
    if t == 'fan': return fan(slide, o)
    if t == 'petri_dish': return petri_dish(slide, o)
    if t == 'material_block': return material_block(slide, o)
    if t == 'dimension_arrow': return dimension_arrow(slide, o)
    if t == 'control_panel': return control_panel(slide, o)
    raise ValueError(f'Unsupported reference-scene object type: {t}')


def render(scene: dict, output: Path) -> Path:
    prs = Presentation()
    cfg = scene.get('slide', {})
    prs.slide_width = Inches(float(cfg.get('width', 10)))
    prs.slide_height = Inches(float(cfg.get('height', 6.84375)))
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = rgb(scene.get('background', '#F7FAFB'))
    for o in scene.get('objects', []): render_object(slide, o)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    check = Presentation(output)
    if not check.slides or not check.slides[0].shapes:
        raise RuntimeError('Reference-scene PPTX verification failed')
    return output


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--scene', required=True, type=Path)
    ap.add_argument('--output', required=True, type=Path)
    a = ap.parse_args()
    scene = json.loads(a.scene.read_text(encoding='utf-8'))
    print(render(scene, a.output.resolve()))


if __name__ == '__main__':
    main()
